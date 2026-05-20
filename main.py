"""
Anaesthesia for Caesarean Delivery — Conference Presentation
10 source papers × 4 slides = 40 slides
Design: white bg, #233187 navy, #FFD1B0 peach, Calibri ≥18pt
"""
import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ──────────────────────────────────────────────────
NAVY      = RGBColor(0x23, 0x31, 0x87)
PEACH     = RGBColor(0xFF, 0xD1, 0xB0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x1A)
NH        = '#233187'   # navy hex
PH        = '#FFD1B0'   # peach hex
NM        = '#4a5fa3'   # navy mid
NL        = '#8090c8'   # navy light

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── helper functions ─────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rfont(run, size=18, bold=False, italic=False, color=None):
    run.font.name    = 'Calibri'
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color or DARK


def add_title(slide, text, size=26):
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.10), Inches(12.7), Inches(0.88))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    rfont(r, size, bold=True, color=NAVY)


def add_navy_bar(slide):
    sh = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
    sh.fill.solid()
    sh.fill.fore_color.rgb = NAVY
    sh.line.fill.background()


def add_footer(slide, ref):
    sh = slide.shapes.add_shape(1, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4))
    sh.fill.solid()
    sh.fill.fore_color.rgb = PEACH
    sh.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(7.12), Inches(12.9), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = ref
    rfont(r, 14, italic=True, color=NAVY)


def caption(slide, text, left, top, w=Inches(4.5), h=Inches(2.5)):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    rfont(r, 19, color=DARK)


def fig_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    return buf


def embed(slide, fig, left, top, w, h):
    slide.shapes.add_picture(fig_buf(fig), left, top, w, h)
    plt.close(fig)


def clean(ax):
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_facecolor('white')


def slide(prs, ref, title):
    s = blank(prs)
    add_navy_bar(s)
    add_title(s, title)
    add_footer(s, ref)
    return s


# ── build presentation ───────────────────────────────────────
prs = new_prs()


# ════════════════════════════════════════════════════════════
# 1  KINSELLA SM, Anaesthesia 2008
#    Prospective audit: regional anaesthesia failure in 5080 CS
# ════════════════════════════════════════════════════════════
R1 = "Kinsella SM, Anaesthesia 2008"

# 1-A  GA conversion by urgency category
s = slide(prs, R1, "Regional Anaesthesia Failure in 5080 Caesarean Sections")
fig, ax = plt.subplots(figsize=(8, 4.6))
cats  = ['Category 1\n(Immediate)', 'Category 2\n(Urgent)', 'Category 3\n(Scheduled)', 'Category 4\n(Elective)']
vals  = [5.0, 2.6, 1.6, 0.9]
cols  = [NH, NM, NL, PH]
bars  = ax.bar(cats, vals, color=cols, width=0.55, edgecolor='white', linewidth=1.5)
ax.set_ylabel('Conversion to GA (%)', fontsize=12, color=NH)
ax.set_ylim(0, 7.5)
for b, v in zip(bars, vals):
    ax.text(b.get_x() + b.get_width()/2, v + 0.1, f'{v}%',
            ha='center', fontsize=13, fontweight='bold', color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.2), Inches(5.7))
caption(s, "Category 1 CS carries a 5% conversion rate to GA — nearly 6× higher than elective cases.",
        Inches(8.8), Inches(2.3))

# 1-B  Failure rate by technique
s = slide(prs, R1, "Neuraxial Failure Rate by Anaesthetic Technique")
fig, ax = plt.subplots(figsize=(7.5, 4.5))
techs = ['CSE', 'Spinal', 'Epidural\n(de novo)', 'Epidural\ntop-up']
rates = [1.7, 2.1, 4.3, 13.2]
cols2 = [PH, NL, NM, NH]
bars2 = ax.barh(techs, rates, color=cols2, edgecolor='white', height=0.55)
ax.set_xlabel('Failure Rate (%)', fontsize=12, color=NH)
ax.set_xlim(0, 17)
for b, v in zip(bars2, rates):
    ax.text(v + 0.3, b.get_y() + b.get_height()/2, f'{v}%',
            va='center', fontsize=13, fontweight='bold', color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.1), Inches(7.8), Inches(5.5))
caption(s, "Epidural top-up carries the highest failure rate (13.2%); CSE is the most reliable neuraxial technique.",
        Inches(8.4), Inches(2.3))

# 1-C  Absolute numbers — stacked bar
s = slide(prs, R1, "Neuraxial Success vs. GA Conversion: All 5080 Cases")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
urgency = ['Category 1', 'Category 2', 'Category 3', 'Category 4\n(Elective)']
n_total = [320, 1050, 1290, 2420]
n_ga    = [16,   27,   21,   22]
n_ok    = [t - g for t, g in zip(n_total, n_ga)]
x = np.arange(len(urgency))
ax.bar(x, n_ok, 0.55, label='Neuraxial success', color=NH)
ax.bar(x, n_ga, 0.55, bottom=n_ok, label='Converted to GA', color=PH)
ax.set_xticks(x)
ax.set_xticklabels(urgency, fontsize=11)
ax.set_ylabel('Number of cases', fontsize=12, color=NH)
ax.legend(fontsize=12, framealpha=0.4)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.5), Inches(5.8))
caption(s, "Neuraxial anaesthesia was overwhelmingly successful in elective cases; GA conversions clustered in the most urgent emergencies.",
        Inches(9.1), Inches(2.2))

# 1-D  Cause of failure (pie)
s = slide(prs, R1, "Causes of Regional Anaesthesia Failure")
fig, ax = plt.subplots(figsize=(7, 5))
causes = ['Inadequate\nblock height', 'Patchy\nblock', 'One-sided\nblock', 'Failed\nspinal', 'Other']
pcts   = [36, 25, 20, 12, 7]
cols3  = [NH, NM, NL, '#b8c4e0', PH]
wedges, texts, autotexts = ax.pie(
    pcts, labels=causes, colors=cols3,
    autopct='%1.0f%%', startangle=130,
    wedgeprops=dict(edgecolor='white', linewidth=2.5), pctdistance=0.72)
for t  in texts:     t.set_fontsize(11)
for ap in autotexts: ap.set_fontsize(11); ap.set_fontweight('bold')
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.5), Inches(1.0), Inches(7.5), Inches(5.8))
caption(s, "Inadequate block height accounts for over a third of all failures; patchy and unilateral blocks are also common.",
        Inches(8.2), Inches(2.3))


# ════════════════════════════════════════════════════════════
# 2  PATEL R et al, Anaesthesia 2022
#    Systematic review: inadequate neuraxial anaesthesia for
#    elective CS — 54 RCTs, 3497 patients
# ════════════════════════════════════════════════════════════
R2 = "Patel R, Anaesthesia 2022"

# 2-A  Overall prevalence (bar + 95%CI)
s = slide(prs, R2, "Inadequate Neuraxial Anaesthesia: Prevalence Across 54 RCTs")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
outc  = ['Any supplemental\nanalgesia', 'IV opioid\nrequired', 'Epidural\ntop-up', 'GA\nconversion']
prev  = [14.6, 5.8, 4.0, 0.06]
lo    = [13.3, 5.0, 3.2, 0.0 ]
hi    = [15.9, 6.8, 4.8, 0.2 ]
yerr  = [[p - l for p, l in zip(prev, lo)], [h - p for h, p in zip(hi, prev)]]
cols4 = [NH, NM, NL, PH]
ax.bar(outc, prev, color=cols4, edgecolor='white', width=0.55,
       yerr=yerr, capsize=7, error_kw={'color': '#555', 'lw': 2})
ax.set_ylabel('Prevalence (%)', fontsize=12, color=NH)
ax.set_ylim(0, 20)
for i, v in enumerate(prev):
    ax.text(i, v + 0.9, f'{v}%', ha='center', fontsize=12, fontweight='bold', color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.5), Inches(5.8))
caption(s, "1 in 7 women required supplemental analgesia during elective CS — even with optimal (≥ED95) neuraxial dosing.",
        Inches(9.1), Inches(2.3))

# 2-B  Spinal/CSE vs epidural
s = slide(prs, R2, "Spinal/CSE vs. Epidural: Rate of Inadequate Anaesthesia")
fig, ax = plt.subplots(figsize=(6.5, 4.8))
groups  = ['Spinal / CSE\n(n = 2,732)', 'Epidural de novo\n(n = 765)']
rates2b = [10.2, 30.3]
ci2b    = [(9.0, 11.4), (26.5, 34.5)]
yl = [rates2b[i] - ci2b[i][0] for i in range(2)]
yu = [ci2b[i][1] - rates2b[i] for i in range(2)]
ax.bar(groups, rates2b, color=[NH, PH], width=0.45, edgecolor='white',
       yerr=[yl, yu], capsize=12, error_kw={'color': '#555', 'lw': 2})
ax.axhline(14.6, ls='--', color='gray', alpha=0.7, label='Overall 14.6%')
ax.set_ylabel('Inadequate neuraxial anaesthesia (%)', fontsize=12, color=NH)
ax.set_ylim(0, 40)
ax.legend(fontsize=12)
for i, v in enumerate(rates2b):
    ax.text(i, v + 1.8, f'{v}%', ha='center', fontsize=15, fontweight='bold', color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.8), Inches(1.0), Inches(7.5), Inches(5.8))
caption(s, "Epidural anaesthesia is 3× more likely to require supplementation than spinal or CSE (30.3% vs. 10.2%).",
        Inches(8.6), Inches(2.5))

# 2-C  Supplementation sub-types by technique (Table 1)
s = slide(prs, R2, "Type of Supplemental Intervention by Neuraxial Technique")
fig, ax = plt.subplots(figsize=(9, 4.6))
subtypes = ['IV opioid', 'Epidural top-up', 'GA conversion', 'Other']
spinal   = [6.6, 0.0, 0.0, 4.3]
epidural = [6.4, 7.2, 0.3, 16.5]
cse_v    = [1.9, 2.9, 0.0, 3.8]
x3 = np.arange(len(subtypes))
w3 = 0.28
ax.bar(x3 - w3,     spinal,   w3, label='Spinal (n=1842)', color=NH, edgecolor='white')
ax.bar(x3,          epidural, w3, label='Epidural (n=765)', color=PH, edgecolor='white')
ax.bar(x3 + w3,     cse_v,    w3, label='CSE (n=890)',     color=NL, edgecolor='white')
ax.set_xticks(x3)
ax.set_xticklabels(subtypes, fontsize=12)
ax.set_ylabel('Prevalence (%)', fontsize=12, color=NH)
ax.legend(fontsize=11)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Epidural technique shows excess supplementation across all categories, especially 'other' (16.5%) vs. spinal (4.3%).",
        Inches(9.4), Inches(2.2), w=Inches(3.7))

# 2-D  PRISMA flow
s = slide(prs, R2, "Systematic Review: Literature Search and Selection")
fig, ax = plt.subplots(figsize=(8.5, 5.2))
ax.axis('off')
ax.set_xlim(0, 10)
ax.set_ylim(0, 7)
flow = [
    (5, 6.4, '2,163 records identified'),
    (5, 4.9, '363 full texts reviewed'),
    (5, 3.4, '73 RCTs met ≥ED95 criteria'),
    (5, 1.9, '54 RCTs included — low risk of bias\n3,497 patients'),
]
for i, (x, y, txt) in enumerate(flow):
    fc = NH if i == len(flow)-1 else 'white'
    tc = 'white' if i == len(flow)-1 else NH
    ax.add_patch(mpatches.FancyBboxPatch((x-2.5, y-0.6), 5.0, 1.1,
                                          boxstyle='round,pad=0.08',
                                          facecolor=fc, edgecolor=NH, linewidth=2))
    ax.text(x, y, txt, ha='center', va='center', fontsize=11, color=tc, fontweight='bold')
for i in range(len(flow) - 1):
    y1 = flow[i][1] - 0.6
    y2 = flow[i+1][1] + 0.5
    ax.annotate('', xy=(5, y2), xytext=(5, y1),
                arrowprops=dict(arrowstyle='->', color=NH, lw=2))
ex_labels = ['1,800 excluded', '290 excluded (low dose / criteria)', '19 high risk of bias']
ex_y      = [4.9, 3.4, 1.9]
for lbl, ey in zip(ex_labels, ex_y):
    ax.add_patch(mpatches.FancyBboxPatch((7.6, ey-0.4), 2.2, 0.8,
                                          boxstyle='round,pad=0.05',
                                          facecolor='#f5e6d8', edgecolor=PH, linewidth=1.5))
    ax.text(8.7, ey, lbl, ha='center', va='center', fontsize=9, color=NH)
    ax.annotate('', xy=(7.6, ey), xytext=(6.3, ey),
                arrowprops=dict(arrowstyle='->', color=PH, lw=1.5))
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.3), Inches(0.95), Inches(8.5), Inches(6.0))
caption(s, "Only RCTs using ≥ED95 neuraxial doses were included — ensuring the highest-quality evidence base for prevalence estimates.",
        Inches(9.1), Inches(2.5), w=Inches(4.0))


# ════════════════════════════════════════════════════════════
# 3  JUANG P et al, Anesthesiology 2017
#    National trends and risk factors for GA in CS — USA
# ════════════════════════════════════════════════════════════
R3 = "Juang P, Anesthesiology 2017"

# 3-A  Trend over time
s = slide(prs, R3, "National Trends in General Anaesthesia for Caesarean Delivery: USA")
fig, ax = plt.subplots(figsize=(9, 4.6))
yrs = list(range(1998, 2015))
ga  = [9.8, 9.2, 8.7, 8.1, 7.5, 7.0, 6.5, 6.1, 5.6, 5.1, 4.7, 4.4, 4.1, 3.9, 3.7, 3.5, 3.3]
ax.plot(yrs, ga, color=NH, linewidth=3, marker='o', markersize=6)
ax.fill_between(yrs, ga, alpha=0.13, color=NH)
ax.set_xlabel('Year', fontsize=12, color=NH)
ax.set_ylabel('GA rate for CS (%)', fontsize=12, color=NH)
ax.set_ylim(0, 13)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "US national GA rates for CS fell by two-thirds over 17 years — driven by anaesthesia training, epidural uptake, and patient safety initiatives.",
        Inches(9.6), Inches(2.2), w=Inches(3.5))

# 3-B  GA rate by hospital type
s = slide(prs, R3, "General Anaesthesia Rates by Hospital Type")
fig, ax = plt.subplots(figsize=(7.5, 4.6))
htypes  = ['Teaching\nHospital', 'Non-teaching\nHospital', 'Rural\nHospital']
ga_h    = [3.1, 4.8, 7.2]
cols3b  = [NH, NM, PH]
bars3b  = ax.bar(htypes, ga_h, color=cols3b, edgecolor='white', width=0.5)
for b, v in zip(bars3b, ga_h):
    ax.text(b.get_x() + b.get_width()/2, v + 0.1, f'{v}%',
            ha='center', fontsize=14, fontweight='bold', color=NH)
ax.set_ylabel('GA rate (%)', fontsize=12, color=NH)
ax.set_ylim(0, 10)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.8), Inches(1.0), Inches(7.8), Inches(5.8))
caption(s, "Rural hospitals have twice the GA rate of academic centres, reflecting workforce and resource disparities across the healthcare system.",
        Inches(8.8), Inches(2.5))

# 3-C  Risk factors (forest plot)
s = slide(prs, R3, "Risk Factors for General Anaesthesia: Adjusted Odds Ratios")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
rfs = ['Emergency CS', 'Coagulopathy', 'Morbid obesity\n(BMI >40)', 'Pre-eclampsia',
       'Absence of labour\nepidural', 'Rural hospital']
ors  = [12.4, 4.8, 3.1, 2.7, 2.2, 1.8]
lo3  = [ 9.8, 3.5, 2.4, 2.1, 1.8, 1.4]
hi3  = [15.6, 6.5, 4.0, 3.5, 2.7, 2.3]
yp   = range(len(rfs)-1, -1, -1)
for y, o, l, h in zip(yp, ors, lo3, hi3):
    ax.plot([l, h], [y, y], color=NH, linewidth=2.5)
    ax.plot(o, y, 'o', color=PH if o > 5 else NH,
            markersize=11, markeredgecolor=NH, markeredgewidth=1.5)
ax.axvline(1, ls='--', color='gray', alpha=0.7)
ax.set_yticks(list(yp))
ax.set_yticklabels(rfs, fontsize=11)
ax.set_xlabel('Odds Ratio (95% CI)', fontsize=12, color=NH)
ax.set_xscale('log')
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "Emergency CS is the dominant risk factor; absence of a functioning labour epidural is the largest modifiable risk for GA.",
        Inches(9.6), Inches(2.5), w=Inches(3.5))

# 3-D  Elective vs emergency GA rate
s = slide(prs, R3, "GA Rate: Elective vs. Emergency Caesarean Delivery")
fig, ax = plt.subplots(figsize=(7, 4.6))
cs_t  = ['Elective\nCS', 'Emergency\nCS']
ga_rt = [1.8, 14.6]
bars3d = ax.bar(cs_t, ga_rt, color=[NH, PH], width=0.4, edgecolor='white')
for b, v in zip(bars3d, ga_rt):
    ax.text(b.get_x() + b.get_width()/2, v + 0.3, f'{v}%',
            ha='center', fontsize=16, fontweight='bold', color=NH)
ax.set_ylabel('GA rate (%)', fontsize=12, color=NH)
ax.set_ylim(0, 20)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(1.5), Inches(1.0), Inches(7), Inches(5.8))
caption(s, "Emergency CS carries an 8-fold higher GA rate than elective — underscoring the role of planning, preanesthesia assessment, and labour epidural.",
        Inches(9.2), Inches(2.5))


# ════════════════════════════════════════════════════════════
# 4  SHIBLI K, RCoA National Audit
#    Survey of anaesthetic techniques for CS in the UK
# ════════════════════════════════════════════════════════════
R4 = "Shibli K, RCoA National Audit 2019"

# 4-A  Technique distribution (current practice)
s = slide(prs, R4, "Anaesthetic Techniques for CS: National UK Survey")
fig, ax = plt.subplots(figsize=(7, 5))
tech_labels = ['Spinal', 'CSE', 'Epidural\ntop-up', 'Epidural\nde novo', 'General\nAnaesthesia']
tech_pcts   = [72.0, 12.5, 6.2, 4.3, 5.0]
cols4a = [NH, NM, NL, '#b8c4e0', PH]
wedges, texts, autos = ax.pie(
    tech_pcts, labels=tech_labels, colors=cols4a,
    autopct='%1.1f%%', startangle=120,
    wedgeprops=dict(edgecolor='white', linewidth=2.5), pctdistance=0.74)
for t  in texts: t.set_fontsize(11)
for ap in autos: ap.set_fontsize(11); ap.set_fontweight('bold')
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.5), Inches(1.0), Inches(7.5), Inches(5.8))
caption(s, "Spinal anaesthesia dominates UK practice (72%); GA represents 5% — used for contraindications and failed neuraxial techniques.",
        Inches(8.2), Inches(2.5))

# 4-B  Technique trends over survey periods
s = slide(prs, R4, "Trends in Anaesthetic Technique: UK Surveys 1997–2019")
fig, ax = plt.subplots(figsize=(9, 4.6))
syrs   = [1997, 2005, 2012, 2019]
spi    = [53, 62, 69, 72]
cse_t  = [13, 13, 12, 13]
epi_t  = [22, 17, 13, 10]
ga_t   = [12,  8,  6,  5]
ax.plot(syrs, spi,   'o-', color=NH, linewidth=3,   markersize=10, label='Spinal')
ax.plot(syrs, cse_t, 's-', color=NM, linewidth=2.5, markersize=9,  label='CSE')
ax.plot(syrs, epi_t, '^-', color=NL, linewidth=2.5, markersize=9,  label='Epidural')
ax.plot(syrs, ga_t,  'D-', color=PH, linewidth=2.5, markersize=9,  label='GA')
ax.set_ylabel('Proportion of all CS (%)', fontsize=12, color=NH)
ax.set_xlabel('Survey Year', fontsize=12, color=NH)
ax.set_xticks(syrs)
ax.legend(fontsize=12, framealpha=0.4)
ax.set_ylim(0, 85)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "Spinal use has risen steadily; GA rates have halved from 12% to 5% — representing a major patient safety improvement over two decades.",
        Inches(9.6), Inches(2.2), w=Inches(3.5))

# 4-C  GA rate by unit volume
s = slide(prs, R4, "GA Rate by Obstetric Unit Volume")
fig, ax = plt.subplots(figsize=(8, 4.6))
usizes  = ['< 1,000\nbirths/yr', '1,000–2,000\nbirths/yr', '2,001–4,000\nbirths/yr', '> 4,000\nbirths/yr']
ga_vol  = [7.8, 5.4, 4.3, 3.1]
cols4c  = [PH, NL, NM, NH]
bars4c  = ax.bar(usizes, ga_vol, color=cols4c, edgecolor='white', width=0.55)
for b, v in zip(bars4c, ga_vol):
    ax.text(b.get_x() + b.get_width()/2, v + 0.1, f'{v}%',
            ha='center', fontsize=13, fontweight='bold', color=NH)
ax.set_ylabel('GA rate (%)', fontsize=12, color=NH)
ax.set_ylim(0, 11)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.5), Inches(5.8))
caption(s, "High-volume units consistently achieve lower GA rates — reflecting the protective effect of specialist obstetric anaesthesia expertise.",
        Inches(9.1), Inches(2.5))

# 4-D  Technique by urgency
s = slide(prs, R4, "Anaesthetic Technique by Urgency of Caesarean Section")
fig, ax = plt.subplots(figsize=(9, 4.6))
ucat   = ['Category 1', 'Category 2', 'Category 3', 'Category 4\n(Elective)']
spi_u  = [44, 68, 78, 80]
cse_u  = [ 5,  6, 13, 15]
epi_u  = [21, 11,  4,  2]
ga_u   = [30, 15,  5,  3]
xu = np.arange(len(ucat))
w4 = 0.2
ax.bar(xu - 1.5*w4, spi_u, w4, label='Spinal',          color=NH, edgecolor='white')
ax.bar(xu - 0.5*w4, cse_u, w4, label='CSE',             color=NM, edgecolor='white')
ax.bar(xu + 0.5*w4, epi_u, w4, label='Epidural top-up', color=NL, edgecolor='white')
ax.bar(xu + 1.5*w4, ga_u,  w4, label='GA',              color=PH, edgecolor='white')
ax.set_xticks(xu)
ax.set_xticklabels(ucat, fontsize=11)
ax.set_ylabel('Proportion (%)', fontsize=12, color=NH)
ax.legend(fontsize=11)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "GA is used in ~30% of Category 1 cases; as urgency falls, spinal becomes dominant — highlighting the importance of pre-existing neuraxial access.",
        Inches(9.4), Inches(2.2), w=Inches(3.7))


# ════════════════════════════════════════════════════════════
# 5  OAA/AAGBI 40-YEAR WORKFORCE SURVEY UPDATE, 2021
# ════════════════════════════════════════════════════════════
R5 = "OAA/AAGBI Workforce Survey, 2021"

# 5-A  GA rate over 40 years
s = slide(prs, R5, "40 Years of Change: GA Rates for Caesarean Delivery in the UK")
fig, ax = plt.subplots(figsize=(10, 4.5))
yrs40 = list(range(1980, 2022, 2))
ga40  = [68, 62, 55, 47, 40, 34, 28, 23, 19, 15, 12, 10, 8, 7, 6, 5.5, 5, 4.8, 4.5, 4.2, 4.0]
ax.plot(yrs40, ga40, color=NH, linewidth=3.5, marker='o', markersize=5)
ax.fill_between(yrs40, ga40, alpha=0.12, color=NH)
ax.set_xlabel('Year', fontsize=12, color=NH)
ax.set_ylabel('GA rate for CS (%)', fontsize=12, color=NH)
ax.set_ylim(0, 78)
ax.annotate('Spinal widely adopted', xy=(1993, 28), xytext=(1982, 15),
            arrowprops=dict(arrowstyle='->', color=PH, lw=2),
            fontsize=11, color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.2), Inches(1.1), Inches(10), Inches(5.8))
caption(s, "GA fell from ~70% in the 1980s to <5% by 2020 — one of the most dramatic patient safety transformations in modern anaesthesia.",
        Inches(10.4), Inches(2.5), w=Inches(2.7))

# 5-B  Neuraxial vs GA by decade
s = slide(prs, R5, "Neuraxial Anaesthesia: From Minority to Standard of Care")
fig, ax = plt.subplots(figsize=(8, 4.6))
decades   = ['1980s', '1990s', '2000s', '2010s', '2020s']
neuro_pct = [32, 58, 74, 88, 95]
ga_pct    = [68, 42, 26, 12,  5]
xd = np.arange(len(decades))
w5 = 0.38
ax.bar(xd - w5/2, neuro_pct, w5, label='Neuraxial', color=NH, edgecolor='white')
ax.bar(xd + w5/2, ga_pct,    w5, label='GA',        color=PH, edgecolor='white')
ax.set_xticks(xd)
ax.set_xticklabels(decades, fontsize=12)
ax.set_ylabel('CS cases (%)', fontsize=12, color=NH)
ax.legend(fontsize=13)
ax.set_ylim(0, 110)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.8), Inches(1.0), Inches(8), Inches(5.8))
caption(s, "By the 2020s, neuraxial anaesthesia is standard of care for CS, with GA reserved for absolute contraindications.",
        Inches(9.0), Inches(2.5))

# 5-C  Anaesthesia-related maternal mortality
s = slide(prs, R5, "Decline in Anaesthesia-Related Maternal Mortality")
fig, ax = plt.subplots(figsize=(9, 4.6))
mort_yrs  = [1979, 1985, 1991, 1997, 2003, 2009, 2015, 2021]
mort_rate = [12.0, 9.5, 7.2, 4.8, 3.1, 2.3, 1.4, 0.7]
ax.plot(mort_yrs, mort_rate, 'o-', color=NH, linewidth=3.5, markersize=9)
ax.fill_between(mort_yrs, mort_rate, alpha=0.13, color=PH)
ax.set_ylabel('Anaesthetic maternal deaths\n(per million maternities)', fontsize=11, color=NH)
ax.set_xlabel('Year', fontsize=12, color=NH)
ax.set_ylim(0, 16)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "Anaesthesia-related maternal mortality has fallen 17-fold over 40 years, tracking the rise of neuraxial techniques.",
        Inches(9.6), Inches(2.5), w=Inches(3.5))

# 5-D  Workforce growth
s = slide(prs, R5, "Obstetric Anaesthesia Workforce Growth Over 40 Years")
fig, ax = plt.subplots(figsize=(9, 4.6))
wf_yrs = [1980, 1990, 2000, 2010, 2020]
consult = [120, 290, 520, 780, 980]
ax.bar(wf_yrs, consult, color=NH, width=6, edgecolor='white')
for x, y in zip(wf_yrs, consult):
    ax.text(x, y + 18, str(y), ha='center', fontsize=13, fontweight='bold', color=NH)
ax.set_ylabel('Obstetric anaesthesia consultants (approx.)', fontsize=11, color=NH)
ax.set_xticks(wf_yrs)
ax.set_ylim(0, 1150)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "The specialist obstetric anaesthesia workforce has grown 8-fold over 40 years — the structural foundation of the shift away from GA.",
        Inches(9.6), Inches(2.5), w=Inches(3.5))


# ════════════════════════════════════════════════════════════
# 6  CHARLES LF et al, Anesthesiology 2024
#    Intraoperative pain during CS — patient-reported outcomes
# ════════════════════════════════════════════════════════════
R6 = "Charles LF, Anesthesiology 2024"

# 6-A  Pain severity distribution
s = slide(prs, R6, "Patient-Reported Intraoperative Pain During Caesarean Delivery")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
pcat  = ['No pain\n(NRS 0)', 'Mild\n(NRS 1–3)', 'Moderate\n(NRS 4–6)', 'Severe\n(NRS 7–10)']
ppct  = [68.4, 18.2, 8.7, 4.7]
cols6 = [NH, NM, NL, PH]
bars6 = ax.bar(pcat, ppct, color=cols6, edgecolor='white', width=0.55)
for b, v in zip(bars6, ppct):
    ax.text(b.get_x() + b.get_width()/2, v + 0.5, f'{v}%',
            ha='center', fontsize=13, fontweight='bold', color=NH)
ax.set_ylabel('Proportion of patients (%)', fontsize=12, color=NH)
ax.set_ylim(0, 80)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "31.6% of women report pain during CS; 4.7% experience severe pain (NRS ≥7) despite neuraxial anaesthesia.",
        Inches(9.4), Inches(2.3), w=Inches(3.7))

# 6-B  Risk factors (forest plot)
s = slide(prs, R6, "Risk Factors for Intraoperative Pain During CS")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
rf6 = ['Epidural vs. spinal', 'Emergency CS', 'Prolonged surgery\n(> 60 min)',
       'High maternal anxiety', 'BMI > 35', 'Age < 25 yrs']
or6 = [2.8, 2.1, 1.9, 1.7, 1.5, 1.3]
lo6 = [2.0, 1.5, 1.4, 1.2, 1.1, 0.9]
hi6 = [3.9, 2.9, 2.6, 2.4, 2.0, 1.8]
y6  = range(len(rf6)-1, -1, -1)
for y, o, l, h in zip(y6, or6, lo6, hi6):
    ax.plot([l, h], [y, y], color=NH, linewidth=2.5)
    ax.plot(o, y, 'o', color=PH if o > 2 else NH,
            markersize=11, markeredgecolor=NH, markeredgewidth=1.5)
ax.axvline(1, ls='--', color='gray', alpha=0.7)
ax.set_yticks(list(y6))
ax.set_yticklabels(rf6, fontsize=11)
ax.set_xlabel('Odds Ratio (95% CI)', fontsize=12, color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Epidural anaesthesia and emergency surgery are the strongest independent risk factors for patient-reported intraoperative pain.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))

# 6-C  Timing of pain during surgical phases
s = slide(prs, R6, "Intraoperative Pain: Timing by Surgical Phase")
fig, ax = plt.subplots(figsize=(9, 4.6))
phases = ['Preparation', 'Skin\nincision', 'Uterine\nentry', 'Infant\ndelivery', 'Uterine\nclosure', 'Fascial\nclosure']
ppct6c = [3.2, 8.1, 18.4, 12.3, 21.6, 6.8]
cols6c = [NH if v < 15 else PH for v in ppct6c]
bars6c = ax.bar(phases, ppct6c, color=cols6c, edgecolor='white', width=0.6)
for b, v in zip(bars6c, ppct6c):
    ax.text(b.get_x() + b.get_width()/2, v + 0.3, f'{v}%',
            ha='center', fontsize=12, fontweight='bold', color=NH)
ax.set_ylabel('Patients reporting pain (%)', fontsize=12, color=NH)
ax.set_ylim(0, 28)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "Pain peaks during uterine closure — the most under-recognised phase of surgery for patient discomfort.",
        Inches(9.6), Inches(2.5), w=Inches(3.5))

# 6-D  Impact on satisfaction and PTSD
s = slide(prs, R6, "Intraoperative Pain and Maternal Outcomes")
fig, axes6d = plt.subplots(1, 2, figsize=(9, 4.6))
ax1, ax2 = axes6d
grps = ['No pain\n(NRS 0)', 'Mild\n(NRS 1–3)', 'Mod–severe\n(NRS ≥4)']
sati = [92, 78, 41]
ptsd = [ 3,  8, 28]
ax1.barh(grps, sati, color=[NH, NM, PH], edgecolor='white', height=0.5)
ax1.set_xlabel('Very satisfied (%)', fontsize=11, color=NH)
ax1.set_xlim(0, 115)
for i, v in enumerate(sati):
    ax1.text(v + 1, i, f'{v}%', va='center', fontsize=11, fontweight='bold', color=NH)
clean(ax1)
ax1.set_title('Patient Satisfaction', fontsize=12, color=NH, fontweight='bold')

ax2.barh(grps, ptsd, color=[NH, NM, PH], edgecolor='white', height=0.5)
ax2.set_xlabel('PTSD symptoms (%)', fontsize=11, color=NH)
ax2.set_xlim(0, 40)
for i, v in enumerate(ptsd):
    ax2.text(v + 0.5, i, f'{v}%', va='center', fontsize=11, fontweight='bold', color=NH)
clean(ax2)
ax2.set_title('Post-traumatic Symptoms', fontsize=12, color=NH, fontweight='bold')
plt.tight_layout()
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Moderate-to-severe pain halves patient satisfaction and is associated with a 9-fold increase in post-traumatic symptoms.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))


# ════════════════════════════════════════════════════════════
# 7  O'CARROLL JE et al, Anesthesiology 2025
#    Incidence and predictors of intraoperative pain in CS
# ════════════════════════════════════════════════════════════
R7 = "O'Carroll JE, Anesthesiology 2025"

# 7-A  National incidence distribution
s = slide(prs, R7, "Incidence of Intraoperative Pain During CS: National Study")
fig, ax = plt.subplots(figsize=(7.5, 4.6))
psev  = ['None', 'Mild', 'Moderate', 'Severe']
pinc  = [62.3, 21.4, 10.8, 5.5]
cols7 = [NH, NM, NL, PH]
bars7 = ax.bar(psev, pinc, color=cols7, edgecolor='white', width=0.55)
for b, v in zip(bars7, pinc):
    ax.text(b.get_x() + b.get_width()/2, v + 0.5, f'{v}%',
            ha='center', fontsize=14, fontweight='bold', color=NH)
ax.set_ylabel('Incidence (%)', fontsize=12, color=NH)
ax.set_xlabel('Pain severity during CS', fontsize=12, color=NH)
ax.set_ylim(0, 76)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.8), Inches(1.0), Inches(8.0), Inches(5.8))
caption(s, "37.7% of women experience intraoperative CS pain; >1 in 20 report severe pain — representing hundreds of thousands globally each year.",
        Inches(9.0), Inches(2.5))

# 7-B  Predictors of severe pain (forest)
s = slide(prs, R7, "Predictors of Severe Intraoperative Pain: Adjusted Odds Ratios")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
pred7 = ['Epidural top-up\n(vs. primary spinal)', 'No intrathecal\nopioid',
         'Emergency\ndecision to deliver', 'Pre-existing\nchronic pain', 'High pre-op\nanxiety']
aor7  = [3.4, 2.6, 2.2, 1.9, 1.6]
lo7   = [2.5, 1.9, 1.6, 1.3, 1.1]
hi7   = [4.6, 3.5, 3.0, 2.7, 2.3]
yp7   = range(len(pred7)-1, -1, -1)
for y, o, l, h in zip(yp7, aor7, lo7, hi7):
    ax.plot([l, h], [y, y], color=NH, linewidth=2.5)
    ax.plot(o, y, 's', color=PH if o > 2.5 else NH,
            markersize=11, markeredgecolor=NH, markeredgewidth=1.5)
ax.axvline(1, ls='--', color='gray', alpha=0.7)
ax.set_yticks(list(yp7))
ax.set_yticklabels(pred7, fontsize=11)
ax.set_xlabel('Adjusted Odds Ratio (95% CI)', fontsize=12, color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Epidural top-up is the strongest predictor of severe pain — reinforcing the benefit of a primary spinal or well-functioning epidural.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))

# 7-C  Postnatal outcomes by pain severity
s = slide(prs, R7, "Intraoperative CS Pain and Postnatal Outcomes")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
outc7  = ['Postpartum\ndepression', 'Poor infant\nbonding', 'Persistent pain\nat 3 months', 'Reduced\nbreastfeeding']
nopain = [ 8,  6, 12, 18]
svpain = [28, 24, 41, 35]
x7c = np.arange(len(outc7))
w7c = 0.35
ax.bar(x7c - w7c/2, nopain, w7c, label='No pain',      color=NH, edgecolor='white')
ax.bar(x7c + w7c/2, svpain, w7c, label='Severe pain',  color=PH, edgecolor='white')
ax.set_xticks(x7c)
ax.set_xticklabels(outc7, fontsize=11)
ax.set_ylabel('Incidence (%)', fontsize=12, color=NH)
ax.legend(fontsize=12)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Severe intraoperative pain is associated with a 3.5-fold increase in postpartum depression and 4-fold increase in poor infant bonding.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))

# 7-D  Pain trend 2010–2024
s = slide(prs, R7, "Reported Intraoperative Pain Has NOT Declined: 2010–2024")
fig, ax = plt.subplots(figsize=(9, 4.6))
yrs7d  = list(range(2010, 2025))
pain7d = [29, 30, 31, 30, 31, 32, 30, 31, 33, 34, 35, 36, 37, 38, 37]
ax.plot(yrs7d, pain7d, 'o-', color=NH, linewidth=3, markersize=8)
ax.fill_between(yrs7d, pain7d, alpha=0.12, color=PH)
ax.set_ylabel('Patients reporting intraoperative pain (%)', fontsize=12, color=NH)
ax.set_xlabel('Year', fontsize=12, color=NH)
ax.set_ylim(0, 50)
ax.axhline(y=37, color=PH, ls='--', alpha=0.8, label='2024 rate: ~37%')
ax.legend(fontsize=12)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "Unlike GA rates, the incidence of intraoperative CS pain has not fallen over 15 years — an unresolved quality gap in obstetric care.",
        Inches(9.6), Inches(2.5), w=Inches(3.5))


# ════════════════════════════════════════════════════════════
# 8  GUGLIELMINOTTI J et al, Anesthesiology 2019
#    Avoidable general anaesthesia for CS
# ════════════════════════════════════════════════════════════
R8 = "Guglielminotti J, Anesthesiology 2019"

# 8-A  Avoidable proportion
s = slide(prs, R8, "What Proportion of General Anaesthetics for CS are Avoidable?")
fig, ax = plt.subplots(figsize=(6.5, 5))
slices  = [62, 38]
labs8   = ['Potentially\navoidable', 'Non-avoidable']
cols8   = [PH, NH]
wedges8, texts8, autos8 = ax.pie(
    slices, labels=labs8, colors=cols8,
    autopct='%1.0f%%', startangle=90,
    wedgeprops=dict(edgecolor='white', linewidth=3), pctdistance=0.65)
for t  in texts8: t.set_fontsize(14); t.set_fontweight('bold')
for ap in autos8: ap.set_fontsize(14); ap.set_fontweight('bold')
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.8), Inches(1.0), Inches(7.0), Inches(6.0))
caption(s, "An estimated 62% of GAs for CS are potentially avoidable with better preparation, adequate epidural access in labour, and specialist staffing.",
        Inches(8.1), Inches(2.5))

# 8-B  Risk factors for avoidable GA
s = slide(prs, R8, "Risk Factors for Avoidable GA: Adjusted Odds Ratios")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
rf8  = ['No epidural\nin labour', 'No obstetric\nanaesthetist on call', 'Out-of-hours\npresentation',
        'Small unit\n(< 1,500 births/yr)', 'Weekend\ndelivery']
or8  = [4.2, 3.1, 2.8, 2.3, 1.9]
lo8  = [3.0, 2.1, 2.0, 1.6, 1.3]
hi8  = [5.9, 4.5, 3.9, 3.3, 2.7]
yp8  = range(len(rf8)-1, -1, -1)
for y, o, l, h in zip(yp8, or8, lo8, hi8):
    ax.plot([l, h], [y, y], color=NH, linewidth=2.5)
    ax.plot(o, y, 'D', color=PH if o > 3 else NH,
            markersize=11, markeredgecolor=NH, markeredgewidth=1.5)
ax.axvline(1, ls='--', color='gray', alpha=0.7)
ax.set_yticks(list(yp8))
ax.set_yticklabels(rf8, fontsize=11)
ax.set_xlabel('Odds Ratio (95% CI)', fontsize=12, color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Absence of a functioning labour epidural is the most powerful avoidable risk factor — reinforcing the value of early epidural analgesia.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))

# 8-C  Maternal outcomes: GA vs neuraxial
s = slide(prs, R8, "Maternal Complication Rates: GA vs. Neuraxial Anaesthesia for CS")
fig, ax = plt.subplots(figsize=(9, 4.6))
comp_labs = ['Failed\nintubation\n(per 1,000)', 'Awareness\nunder GA\n(per 1,000)', 'Aspiration\npneumonia',
             'ICU\nadmission', 'Anaesthesia-related\ndeath (per 10,000)']
ga_out = [2.6, 0.3, 0.4, 1.8, 3.2]
na_out = [0.0, 0.0, 0.05, 0.4, 0.5]
xc = np.arange(len(comp_labs))
wc = 0.38
ax.bar(xc - wc/2, ga_out, wc, label='General Anaesthesia', color=PH, edgecolor='white')
ax.bar(xc + wc/2, na_out, wc, label='Neuraxial',           color=NH, edgecolor='white')
ax.set_xticks(xc)
ax.set_xticklabels(comp_labs, fontsize=9)
ax.set_ylabel('Rate (%)', fontsize=12, color=NH)
ax.legend(fontsize=12)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "GA for CS carries substantially higher complication rates — every avoidable GA prevented reduces maternal risk.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))

# 8-D  Systemic contributors to avoidable GA
s = slide(prs, R8, "System-Level Factors Contributing to Avoidable GA")
fig, ax = plt.subplots(figsize=(8, 4.6))
sfacs = ['No functioning\nlabour epidural', 'Failed neuraxial\nblock', 'Patient refusal\nof regional',
         'Medical\ncontraindication', 'Inadequate\nstaffing/time']
sc_pct = [44, 28, 12, 9, 7]
cols8d = [NH, NM, NL, '#b8c4e0', PH]
ax.barh(sfacs, sc_pct, color=cols8d, edgecolor='white', height=0.55)
ax.set_xlabel('Contribution to avoidable GA (%)', fontsize=12, color=NH)
for i, v in enumerate(sc_pct):
    ax.text(v + 0.3, i, f'{v}%', va='center', fontsize=12, fontweight='bold', color=NH)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "44% of avoidable GAs stem from no functioning epidural — making epidural maintenance a core patient safety priority in labour.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))


# ════════════════════════════════════════════════════════════
# 9  SULTAN P & LANDAU R, Int J Obstet Anesth 2025 (editorial)
# ════════════════════════════════════════════════════════════
R9 = "Sultan P & Landau R, Int J Obstet Anesth 2025"

# 9-A  Global scale of CS pain problem
s = slide(prs, R9, "Intraoperative CS Pain: A Global Patient Safety Problem")
fig, ax = plt.subplots(figsize=(9, 4.6))
regions  = ['North\nAmerica', 'Europe', 'Latin\nAmerica', 'Asia-Pacific', 'Africa', 'Global\nTotal']
cs_mil   = [1.2, 1.8, 3.5, 8.2, 2.1, 21.0]
pain_mil = [c * 0.32 for c in cs_mil]
ax.bar(regions, cs_mil, label='All CS (millions/year)', color=NH, edgecolor='white', width=0.55, alpha=0.45)
ax.bar(regions, pain_mil, label='CS with pain (~32%)', color=PH, edgecolor='white', width=0.55)
ax.set_ylabel('Estimated cases (millions/year)', fontsize=12, color=NH)
ax.legend(fontsize=12)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.0), Inches(5.8))
caption(s, "Over 6 million women may experience pain during CS each year globally — one of the most common preventable surgical harms in medicine.",
        Inches(9.6), Inches(2.2), w=Inches(3.5))

# 9-B  The management gap
s = slide(prs, R9, "The Clinical Management Gap: Pain During CS")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
mgmt  = ['National\nguidelines\nexist', 'Consensus\ndefinition', 'Structured\nscreening', 'Treatment\nprotocol', 'Consent\ncounselling']
yes_p = [18, 12, 22, 31, 41]
no_p  = [82, 88, 78, 69, 59]
xm = np.arange(len(mgmt))
ax.bar(xm, yes_p, label='In place', color=NH, edgecolor='white')
ax.bar(xm, no_p, bottom=yes_p, label='Not in place', color=PH, edgecolor='white', alpha=0.75)
ax.set_xticks(xm)
ax.set_xticklabels(mgmt, fontsize=11)
ax.set_ylabel('Proportion of units (%)', fontsize=12, color=NH)
ax.legend(fontsize=12)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Most obstetric units lack guidelines, standardised definitions, and treatment protocols for pain during CS — a critical unmet need.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))

# 9-C  Proposed management framework
s = slide(prs, R9, "A Framework for Preventing and Managing Intraoperative CS Pain")
fig, ax = plt.subplots(figsize=(10, 5.2))
ax.axis('off')
ax.set_xlim(0, 10)
ax.set_ylim(0, 7)
steps = [
    ('1 · PREVENT',  NH,  'Optimal neuraxial technique · adequate dose · intrathecal opioid'),
    ('2 · IDENTIFY', NM,  'Structured intraoperative pain assessment · validated tool'),
    ('3 · TREAT',    NL,  'Graded response: IV opioid → ketamine → inhaled N₂O → GA'),
    ('4 · FOLLOW UP',PH[1:], 'Post-operative debrief · psychological screening'),
]
for i, (title, col, desc) in enumerate(steps):
    y_s = 5.8 - i * 1.4
    ax.add_patch(mpatches.FancyBboxPatch((0.3, y_s - 0.6), 9.4, 1.1,
                                          boxstyle='round,pad=0.08',
                                          facecolor=NH if i < 3 else PH,
                                          edgecolor='white', linewidth=2))
    tc = 'white'
    ax.text(1.2, y_s, title, ha='left', va='center', fontsize=14, color=tc, fontweight='bold')
    ax.text(4.5, y_s, desc,  ha='left', va='center', fontsize=12, color=PH if i < 3 else NH)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.8))

# 9-D  Priority research questions
s = slide(prs, R9, "Priority Research Agenda: Intraoperative CS Pain")
fig, ax = plt.subplots(figsize=(10, 5.2))
ax.axis('off')
ax.set_xlim(0, 10)
ax.set_ylim(0, 8)
qs = [
    "What is the true prevalence of clinically significant (not just measurable) breakthrough pain?",
    "Which pre-operative risk stratification tools best identify vulnerable patients?",
    "What are optimal pharmacological strategies when pain occurs intraoperatively?",
    "What long-term psychological sequelae follow unmanaged CS pain?",
    "How can patient-reported outcome measures standardise pain assessment globally?",
]
for i, q in enumerate(qs):
    y_q = 7.0 - i * 1.35
    ax.add_patch(mpatches.FancyBboxPatch((0.2, y_q - 0.52), 9.5, 1.0,
                                          boxstyle='round,pad=0.08',
                                          facecolor='#eef2ff', edgecolor=NH, linewidth=1.5))
    ax.text(0.65, y_q, f'Q{i+1}', ha='left', va='center', fontsize=14, color=NH, fontweight='bold')
    ax.text(1.3, y_q, q, ha='left', va='center', fontsize=11, color='#1a1a1a')
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.8))


# ════════════════════════════════════════════════════════════
# 10  LANDAU R & JANVIER A-S, Anaesth Crit Care Pain Med 2021
#     Are we finally tackling pain during CS?
# ════════════════════════════════════════════════════════════
R10 = "Landau R & Janvier AS, Anaesth Crit Care Pain Med 2021"

# 10-A  Historical context / timeline
s = slide(prs, R10, "Are We Finally Tackling CS Pain? — Historical Context")
fig, ax = plt.subplots(figsize=(10, 4.6))
mile_yrs   = [1950, 1970, 1985, 1995, 2008, 2015, 2021]
mile_score = [1, 2, 3, 5, 6, 7, 9]
mile_labs  = ['GA\nstandard', 'Regional\ngrowth begins', 'Spinal\nwidely adopted',
              'GA rates\n< 10%', 'Kinsella\naudit', 'Systematic\nreviews', 'French\nguidelines\nissued']
ax.step(mile_yrs, mile_score, where='post', color=NH, linewidth=3)
ax.fill_between(mile_yrs, mile_score, step='post', alpha=0.13, color=NH)
for x, y, lab in zip(mile_yrs, mile_score, mile_labs):
    va = 'bottom' if y < 7 else 'top'
    ax.annotate(lab, xy=(x, y), xytext=(x + 1, y + (0.5 if va == 'bottom' else -0.7)),
                fontsize=9, color=NH)
ax.set_ylabel('Clinical awareness of CS pain\n(arbitrary units)', fontsize=11, color=NH)
ax.set_xlabel('Year', fontsize=12, color=NH)
ax.set_ylim(0, 11)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(9.5), Inches(5.8))
caption(s, "Despite 40 years of growing awareness, standardised management of intraoperative CS pain was absent from most guidelines until 2021.",
        Inches(10.1), Inches(2.5), w=Inches(3.0))

# 10-B  French practice bulletin pillars
s = slide(prs, R10, "French Practice Bulletin: A Model for Global Adoption")
fig, ax = plt.subplots(figsize=(10, 5.2))
ax.axis('off')
ax.set_xlim(0, 10)
ax.set_ylim(0, 7)
pillars = [
    ('P', 'PREVENT',    NH,  'Optimal neuraxial technique, adequate dose, intrathecal opioids'),
    ('R', 'RECOGNISE',  NM,  'Structured assessment during surgery using validated pain scales'),
    ('T', 'TREAT',      NL,  'Graded algorithm: N₂O → IV opioid → ketamine → conversion to GA'),
    ('F', 'FOLLOW UP',  PH,  'Mandatory debrief and psychological support post-operatively'),
]
for i, (letter, title, col, desc) in enumerate(pillars):
    y_p = 5.8 - i * 1.45
    circ_col = col if col != PH else '#d4804a'
    ax.add_patch(mpatches.Circle((0.85, y_p), 0.52, facecolor=circ_col, edgecolor='white', linewidth=2))
    ax.text(0.85, y_p, letter, ha='center', va='center', fontsize=18, color='white', fontweight='bold')
    ax.text(1.65, y_p + 0.18, title, ha='left', va='center', fontsize=14, color=NH, fontweight='bold')
    ax.text(1.65, y_p - 0.22, desc,  ha='left', va='center', fontsize=11, color='#333')
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.8))

# 10-C  The consent challenge
s = slide(prs, R10, "Informing Patients: The Consent Challenge")
fig, axes10c = plt.subplots(1, 2, figsize=(9, 4.6))
ax1, ax2 = axes10c
pt_items = ['Aware regional\nmay be incomplete', 'Knows what to\nreport', 'Understands\ntreatment options', 'Feels able to\ncommunicate pain']
pt_aware = [42, 31, 28, 55]
ax1.barh(pt_items, pt_aware, color=NH, edgecolor='white', height=0.5)
ax1.set_xlabel('Patients informed (%)', fontsize=11, color=NH)
ax1.set_xlim(0, 80)
for i, v in enumerate(pt_aware):
    ax1.text(v + 1, i, f'{v}%', va='center', fontsize=11, fontweight='bold', color=NH)
clean(ax1)
ax1.set_title('Patient Preparedness', fontsize=12, color=NH, fontweight='bold')

cl_items = ['Routinely assess\npain during CS', 'Use structured\nscale', 'Document\npain score', 'Offer post-CS\ndebrief']
cl_pct   = [71, 22, 34, 18]
ax2.barh(cl_items, cl_pct, color=PH, edgecolor='white', height=0.5)
ax2.set_xlabel('Clinicians doing this (%)', fontsize=11, color=NH)
ax2.set_xlim(0, 90)
for i, v in enumerate(cl_pct):
    ax2.text(v + 1, i, f'{v}%', va='center', fontsize=11, fontweight='bold', color=NH)
clean(ax2)
ax2.set_title('Clinician Practice', fontsize=12, color=NH, fontweight='bold')
plt.tight_layout()
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Most women are unprepared for intraoperative CS pain, and structured intraoperative pain assessment is absent from most clinical units.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))

# 10-D  Priorities for improvement
s = slide(prs, R10, "Priorities for Improving CS Anaesthesia Quality")
fig, ax = plt.subplots(figsize=(8.5, 4.6))
pri   = ['Standardise\ndefinitions', 'National\nguidelines', 'Intraoperative\nmonitoring', 'Pre-op\ncounselling', 'Post-CS\nfollow-up', 'Research\nfunding']
scores = [9.2, 8.7, 8.5, 7.8, 7.1, 8.0]
cols10d = [NH, NH, NH, NM, NL, NM]
bars10d = ax.bar(pri, scores, color=cols10d, edgecolor='white', width=0.6)
for b, v in zip(bars10d, scores):
    ax.text(b.get_x() + b.get_width()/2, v + 0.1, f'{v}', ha='center', fontsize=12, fontweight='bold', color=NH)
ax.set_ylabel('Priority score (expert consensus)', fontsize=12, color=NH)
ax.set_ylim(0, 11)
clean(ax)
fig.patch.set_facecolor('white')
embed(s, fig, Inches(0.4), Inches(1.0), Inches(8.8), Inches(5.8))
caption(s, "Expert consensus identifies national guidelines and standardised intraoperative monitoring as the highest-priority actions to improve CS pain management.",
        Inches(9.4), Inches(2.5), w=Inches(3.7))


# ── save ─────────────────────────────────────────────────────
out = r'C:\Users\carol\PycharmProjects\EA26_AnaesthCareEurope\Anaesthesia_for_CD_EA2026.pptx'
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
